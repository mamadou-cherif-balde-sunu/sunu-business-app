# relation_client/pptx_report.py
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from django.conf import settings
from pptx.util import Pt
from .utils_report import (
    interpretation_nps, interpretation_csat, interpretation_ces, conclusion_perimetre
)

TEMPLATE_PATH = os.path.join(settings.BASE_DIR, "relation_client", "templates", "relation_client", "Template_Rapport_SUNU.pptx")


def _set_text(shape, text):
    """Remplace tout le texte d'une zone en gardant le formatage du premier run."""
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = text
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = text
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _shape_by_id(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _update_pie_chart(shape, values):
    """values: liste de % (0-100) dans l'ordre des catégories déjà présentes."""
    chart = shape.chart
    cats = list(chart.plots[0].categories)
    data = CategoryChartData()
    data.categories = cats
    data.add_series(chart.series[0].name or "Série", [v / 100 for v in values])
    chart.replace_data(data)


def _safe_pct(n, d):
    return f"{round(n/d*100):.0f}%" if d else "0%"


def _set_cell_small(cell, text, size=8):
    """Écrit le texte dans la cellule en forçant une petite police,
    pour éviter que la ligne grossisse et fasse chevaucher les tableaux suivants."""
    cell.text = text
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(size)
        for run in para.runs:
            run.font.size = Pt(size)


def _set_table_header(tbl, period_label):
    """Remplace l'en-tête 'Sem. du 08 au 12 Déc. 25' par la vraie période, en
    version courte et petite police pour ne pas agrandir la ligne (ce qui
    ferait chevaucher les tableaux suivants sur la slide)."""
    short_label = (period_label
                   .replace("Semaine du ", "")
                   .replace("Semaine ", ""))
    _set_cell_small(tbl.cell(0, 2), short_label, size=8)
    _set_cell_small(tbl.cell(0, 6), f"Total {short_label}", size=8)
    _set_cell_small(tbl.cell(0, 8), "", size=8)


def _clear_qr_appel_cols(tbl, row_start, row_end):
    """Vide (au lieu de mettre '-') les colonnes QR Code / Appels et Var. Nbre :
    non gérées par l'application, à remplir manuellement par les utilisateurs."""
    for r in range(row_start, row_end):
        for c in (2, 3, 4, 5, 9):
            _set_cell_small(tbl.cell(r, c), "", size=8)


def generate_report(period_label: str, stats: dict, output_path: str):
    """
    period_label : ex "Semaine du 08 au 12 Déc. 2025" ou "Janvier 2026"
    stats : dict retourné par get_period_stats() -> {"vie":..., "iard":..., "global":...}
    output_path : chemin complet du .pptx à générer
    """
    prs = Presentation(TEMPLATE_PATH)
    vie, iard, glob = stats["vie"], stats["iard"], stats["global"]

    # ---- SLIDE 1 : page de titre ----
    s1 = prs.slides[0]
    _set_text(_shape_by_id(s1, 2), f"RESTITUTION SONDAGE \nNPS - CSAT - CES \n{period_label}")
    _set_text(_shape_by_id(s1, 6), f"Département Relation Clients\nSUNU Business GIE Côte d'Ivoire\n{period_label}")

    # ---- SLIDE 4 : méthodologie (ciblage, échantillon, période) ----
    s4 = prs.slides[3]
    shp4_3 = _shape_by_id(s4, 3)
    if shp4_3:
        _set_text(shp4_3,
                  f" Etape 1: Le cadrage\nCiblage : Les clients ayant contacté le Call Center ou "
                  f"s'étant rendu aux Espaces Conseils sur la période {period_label} "
                  f"(périmètre VIE et IARD)\n\n"
                  f"Echantillon :\n"
                  f"{vie['total']} clients pour la VIE\n"
                  f"{iard['total']} clients pour l'IARD\n"
                  f"Total base échantillon : {glob['total']} clients")
    shp4_4 = _shape_by_id(s4, 4)
    if shp4_4:
        _set_text(shp4_4,
                  f"Etape 2: L'élaboration du questionnaire\n"
                  f"Il est composé principalement de questions reparties selon les thèmes suivants : "
                  f"le Customer Satisfaction (CSAT), le Customer Effort Score (CES) et le Net "
                  f"Promotion Score (NPS)\n\n"
                  f"Etape 3 : L'administration de l'enquête\n"
                  f"L'enquête se fait via QR Code et Appels\n"
                  f"Période : {period_label}")

    # ---- SLIDE 6 : présentation globale chiffres ----
    s6 = prs.slides[5]
    _set_text(_shape_by_id(s6, 12),
              f"Echantillon : \n {vie['total']} clients ont répondu pour le périmètre VIE\n"
              f" {iard['total']} clients ont répondu pour le périmètre IARD\n"
              f"Total base échantillon {period_label} : {glob['total']} clients")

    # --- Tableau NPS (id=17) ---
    tbl_nps = _shape_by_id(s6, 17).table
    _set_table_header(tbl_nps, period_label)
    _set_cell_small(tbl_nps.cell(3, 6), str(vie["promoteurs"]));   _set_cell_small(tbl_nps.cell(3, 7), _safe_pct(vie["promoteurs"], vie["total"]))
    _set_cell_small(tbl_nps.cell(4, 6), str(vie["passifs"]));      _set_cell_small(tbl_nps.cell(4, 7), _safe_pct(vie["passifs"], vie["total"]))
    _set_cell_small(tbl_nps.cell(5, 6), str(vie["detracteurs"]));  _set_cell_small(tbl_nps.cell(5, 7), _safe_pct(vie["detracteurs"], vie["total"]))
    _set_cell_small(tbl_nps.cell(6, 6), str(vie["total"]));        _set_cell_small(tbl_nps.cell(6, 7), f"{vie['nps']}%")
    _set_cell_small(tbl_nps.cell(7, 6), str(iard["promoteurs"]));  _set_cell_small(tbl_nps.cell(7, 7), _safe_pct(iard["promoteurs"], iard["total"]))
    _set_cell_small(tbl_nps.cell(8, 6), str(iard["passifs"]));     _set_cell_small(tbl_nps.cell(8, 7), _safe_pct(iard["passifs"], iard["total"]))
    _set_cell_small(tbl_nps.cell(9, 6), str(iard["detracteurs"])); _set_cell_small(tbl_nps.cell(9, 7), _safe_pct(iard["detracteurs"], iard["total"]))
    _set_cell_small(tbl_nps.cell(10, 6), str(iard["total"]));      _set_cell_small(tbl_nps.cell(10, 7), f"{iard['nps']}%")
    _set_cell_small(tbl_nps.cell(11, 6), str(glob["total"]))
    _set_cell_small(tbl_nps.cell(12, 6), f"{glob['nps']}%")
    _clear_qr_appel_cols(tbl_nps, 3, 13)

    # --- Tableau CSAT (id=18) ---
    tbl_csat = _shape_by_id(s6, 18).table
    _set_table_header(tbl_csat, period_label)
    _set_cell_small(tbl_csat.cell(3, 6), str(vie["satisfaits"]));    _set_cell_small(tbl_csat.cell(3, 7), _safe_pct(vie["satisfaits"], vie["total"]))
    _set_cell_small(tbl_csat.cell(4, 6), str(vie["insatisfaits"]));  _set_cell_small(tbl_csat.cell(4, 7), _safe_pct(vie["insatisfaits"], vie["total"]))
    _set_cell_small(tbl_csat.cell(5, 6), str(vie["total"]));         _set_cell_small(tbl_csat.cell(5, 7), f"{vie['csat']}%")
    _set_cell_small(tbl_csat.cell(6, 6), str(iard["satisfaits"]));   _set_cell_small(tbl_csat.cell(6, 7), _safe_pct(iard["satisfaits"], iard["total"]))
    _set_cell_small(tbl_csat.cell(7, 6), str(iard["insatisfaits"])); _set_cell_small(tbl_csat.cell(7, 7), _safe_pct(iard["insatisfaits"], iard["total"]))
    _set_cell_small(tbl_csat.cell(8, 6), str(iard["total"]));        _set_cell_small(tbl_csat.cell(8, 7), f"{iard['csat']}%")
    _set_cell_small(tbl_csat.cell(9, 6), str(glob["total"]))
    _set_cell_small(tbl_csat.cell(10, 6), f"{glob['csat']}%")
    _clear_qr_appel_cols(tbl_csat, 3, 11)

    # --- Tableau CES (id=19) ---
    tbl_ces = _shape_by_id(s6, 19).table
    _set_table_header(tbl_ces, period_label)
    _set_cell_small(tbl_ces.cell(3, 6), str(vie["effort_faible"]));   _set_cell_small(tbl_ces.cell(3, 7), _safe_pct(vie["effort_faible"], vie["total"]))
    _set_cell_small(tbl_ces.cell(4, 6), str(vie["effort_modere"]));   _set_cell_small(tbl_ces.cell(4, 7), _safe_pct(vie["effort_modere"], vie["total"]))
    _set_cell_small(tbl_ces.cell(5, 6), str(vie["effort_eleve"]));    _set_cell_small(tbl_ces.cell(5, 7), _safe_pct(vie["effort_eleve"], vie["total"]))
    _set_cell_small(tbl_ces.cell(6, 6), str(vie["total"]));           _set_cell_small(tbl_ces.cell(6, 7), f"{vie['ces']}%")
    _set_cell_small(tbl_ces.cell(7, 6), str(iard["effort_faible"]));  _set_cell_small(tbl_ces.cell(7, 7), _safe_pct(iard["effort_faible"], iard["total"]))
    _set_cell_small(tbl_ces.cell(8, 6), str(iard["effort_modere"]));  _set_cell_small(tbl_ces.cell(8, 7), _safe_pct(iard["effort_modere"], iard["total"]))
    _set_cell_small(tbl_ces.cell(9, 6), str(iard["effort_eleve"]));   _set_cell_small(tbl_ces.cell(9, 7), _safe_pct(iard["effort_eleve"], iard["total"]))
    _set_cell_small(tbl_ces.cell(10, 6), str(iard["total"]));         _set_cell_small(tbl_ces.cell(10, 7), f"{iard['ces']}%")
    _set_cell_small(tbl_ces.cell(11, 6), str(glob["total"]))
    _set_cell_small(tbl_ces.cell(12, 6), f"{glob['ces']}%")
    _clear_qr_appel_cols(tbl_ces, 3, 13)

    # Pastilles rondes : UNIQUEMENT NPS/CSAT/CES, tout le reste vidé (Appels, QR Code, Echantillon, sous-textes)
    _set_cell_small(_shape_by_id(s6, 6).table.cell(0, 0), f"NPS : {glob['nps']}%", size=14)
    _set_cell_small(_shape_by_id(s6, 7).table.cell(0, 0), f"CSAT : {glob['csat']}%", size=14)
    _set_cell_small(_shape_by_id(s6, 8).table.cell(0, 0), f"CES : {glob['ces']}%", size=14)

    # ---- SLIDE 7 : périmètre VIE ----
    s7 = prs.slides[6]
    _set_text(_shape_by_id(s7, 3),
              f"{vie['csat']}% des clients interrogés sont satisfaits du traitement de leurs "
              f"requêtes contre {100 - vie['csat']:.1f}% d'insatisfaits")
    _shape_by_id(s7, 5).table.cell(0, 0).text = "Total"
    _shape_by_id(s7, 5).table.cell(0, 1).text = str(vie["total"])
    _shape_by_id(s7, 5).table.cell(0, 2).text = "100%"
    _set_text(_shape_by_id(s7, 6),
              f"{round(vie['effort_faible']/vie['total']*100,1) if vie['total'] else 0}% des clients "
              f"interrogés estiment avoir fourni un faible effort pour avoir accès à l'information "
              f"recherchée au niveau de la VIE")
    _set_text(_shape_by_id(s7, 7),
              f"A la question de savoir, quelle est la probabilité pour vous de recommander SUNU "
              f"Assurances comme assureur à vos proches, amis et connaissances ? "
              f"{interpretation_nps(vie['nps'])}")
    _set_text(_shape_by_id(s7, 8), "Conclusion :\n" + conclusion_perimetre("VIE", vie))

    total_v = vie["total"] or 1
    _update_pie_chart(_shape_by_id(s7, 9), [
        vie["detracteurs"] / total_v * 100, vie["passifs"] / total_v * 100, vie["promoteurs"] / total_v * 100
    ])
    _update_pie_chart(_shape_by_id(s7, 10), [
        vie["effort_faible"] / total_v * 100, vie["effort_modere"] / total_v * 100, vie["effort_eleve"] / total_v * 100
    ])
    _update_pie_chart(_shape_by_id(s7, 13), [
        vie["satisfaits"] / total_v * 100, vie["insatisfaits"] / total_v * 100
    ])

    # ---- SLIDE 8 : périmètre IARD ----
    s8 = prs.slides[7]
    _set_text(_shape_by_id(s8, 3),
              f"{iard['csat']}% des clients interrogés sont satisfaits du traitement de leurs "
              f"requêtes contre {100 - iard['csat']:.1f}% d'insatisfaits")
    _shape_by_id(s8, 5).table.cell(0, 0).text = "Total"
    _shape_by_id(s8, 5).table.cell(0, 1).text = str(iard["total"])
    _shape_by_id(s8, 5).table.cell(0, 2).text = "100%"
    _set_text(_shape_by_id(s8, 6),
              f"{round(iard['effort_faible']/iard['total']*100,1) if iard['total'] else 0}% des clients "
              f"interrogés estiment avoir fourni un faible effort pour avoir accès à l'information "
              f"recherchée au niveau de l'IARD")
    _set_text(_shape_by_id(s8, 7),
              f"A la question de savoir, quelle est la probabilité pour vous de recommander SUNU "
              f"Assurances comme assureur à vos proches, amis et connaissances ? "
              f"{interpretation_nps(iard['nps'])}")
    _set_text(_shape_by_id(s8, 8), "Conclusion :\n" + conclusion_perimetre("IARD", iard))

    total_i = iard["total"] or 1
    _update_pie_chart(_shape_by_id(s8, 11), [
        iard["detracteurs"] / total_i * 100, iard["passifs"] / total_i * 100, iard["promoteurs"] / total_i * 100
    ])
    _update_pie_chart(_shape_by_id(s8, 13), [
        iard["effort_faible"] / total_i * 100, iard["effort_modere"] / total_i * 100, iard["effort_eleve"] / total_i * 100
    ])
    _update_pie_chart(_shape_by_id(s8, 14), [
        iard["satisfaits"] / total_i * 100, iard["insatisfaits"] / total_i * 100
    ])

    prs.save(output_path)
    return output_path